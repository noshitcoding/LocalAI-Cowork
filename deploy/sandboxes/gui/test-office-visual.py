#!/usr/bin/env python3
import json
import subprocess
import sys
from pathlib import Path

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches, Pt as PptxPt


WORKSPACE = Path("/workspace")
TOOL = Path("/opt/cowork/office-tool.py")
MARKERS = {
    "docx": "Open Cowork Writer visual baseline",
    "xlsx": "Open Cowork Calc visual baseline",
    "pptx": "Open Cowork Impress visual baseline",
}


def run_tool(payload: dict) -> dict:
    result = subprocess.run(
        [sys.executable, str(TOOL)],
        input=json.dumps(payload),
        text=True,
        capture_output=True,
        timeout=300,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(f"office-tool failed for {payload['action']}: {result.stderr[-4000:]}")
    return json.loads(result.stdout)


def set_cell_fill(cell, color: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), color)
    properties.append(shading)


def create_docx(path: Path) -> None:
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(MARKERS["docx"])
    run.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0x16, 0x6B, 0x8F)
    document.add_paragraph(
        "This page verifies deterministic OOXML structure, LibreOffice PDF export, and PNG review rendering."
    )
    table = document.add_table(rows=3, cols=2)
    table.style = "Table Grid"
    values = (("Capability", "Status"), ("Server runtime", "Ready"), ("Visual review", "Passed"))
    for row_index, values_row in enumerate(values):
        for column_index, value in enumerate(values_row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            if row_index == 0:
                set_cell_fill(cell, "166B8F")
                for paragraph in cell.paragraphs:
                    for item in paragraph.runs:
                        item.bold = True
                        item.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    document.add_paragraph("Visual anchor: COWORK-DOCX-2026")
    document.save(path)


def create_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Visual acceptance"
    sheet.append([MARKERS["xlsx"], "Value"])
    sheet.append(["Linux GUI sandbox", 42])
    sheet.append(["LibreOffice export", 67])
    sheet.append(["Review artifact", 91])
    header_fill = PatternFill("solid", fgColor="166B8F")
    for cell in sheet[1]:
        cell.fill = header_fill
        cell.font = Font(color="FFFFFF", bold=True, size=14)
        cell.alignment = Alignment(horizontal="center")
    sheet.column_dimensions["A"].width = 40
    sheet.column_dimensions["B"].width = 18
    chart = BarChart()
    chart.title = "Distributed runtime readiness"
    chart.y_axis.title = "Score"
    chart.x_axis.title = "Area"
    chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=4), titles_from_data=True)
    chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=4))
    chart.height = 8
    chart.width = 14
    sheet.add_chart(chart, "D2")
    sheet.sheet_properties.pageSetUpPr.fitToPage = True
    sheet.page_setup.fitToWidth = 1
    sheet.page_setup.fitToHeight = 1
    sheet.page_setup.orientation = "landscape"
    sheet.print_area = "A1:K20"
    workbook.save(path)


def create_pptx(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333)
    presentation.slide_height = PptxInches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = PptxRGBColor(0x12, 0x1B, 0x2A)
    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PptxInches(0.8),
        PptxInches(0.8),
        PptxInches(11.7),
        PptxInches(1.3),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = PptxRGBColor(0x16, 0xA3, 0xB6)
    banner.line.fill.background()
    title = banner.text_frame.paragraphs[0]
    title.text = MARKERS["pptx"]
    title.alignment = PP_ALIGN.CENTER
    title.runs[0].font.bold = True
    title.runs[0].font.size = PptxPt(28)
    title.runs[0].font.color.rgb = PptxRGBColor(0xFF, 0xFF, 0xFF)
    body = slide.shapes.add_textbox(
        PptxInches(1.2), PptxInches(2.8), PptxInches(10.9), PptxInches(2.5)
    ).text_frame
    body.text = "Linux GUI sandbox\nLibreOffice rendering\nVersioned review artifacts"
    for paragraph in body.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = PptxPt(24)
            run.font.color.rgb = PptxRGBColor(0xF3, 0xF4, 0xF6)
    presentation.save(path)


def image_metrics(path: Path) -> dict:
    with Image.open(path) as source:
        image = source.convert("RGB")
        image.thumbnail((800, 800))
        pixels = list(image.getdata())
        non_white = sum(1 for red, green, blue in pixels if min(red, green, blue) < 245)
        saturated = sum(1 for red, green, blue in pixels if max(red, green, blue) - min(red, green, blue) > 35)
        dark = sum(1 for red, green, blue in pixels if red + green + blue < 240)
        return {
            "width": source.width,
            "height": source.height,
            "non_white_ratio": non_white / len(pixels),
            "saturated_ratio": saturated / len(pixels),
            "dark_ratio": dark / len(pixels),
        }


def verify_render(kind: str, path: Path) -> dict:
    inspected = run_tool({"action": "inspect", "path": path.name})
    if inspected["type"] not in {"word", "excel", "powerpoint"}:
        raise AssertionError(f"unexpected structural inspection for {kind}: {inspected}")
    preview = run_tool({"action": "preview", "path": path.name, "all_pages": True, "dpi": 120})
    images = preview.get("images") or []
    if len(images) != 1:
        raise AssertionError(f"{kind} visual fixture rendered {len(images)} pages instead of one")
    pdf_inspection = run_tool({"action": "inspect", "path": preview["pdf"]})
    if pdf_inspection.get("pages") != 1 or MARKERS[kind] not in pdf_inspection.get("text", ""):
        raise AssertionError(f"{kind} PDF lost its page count or semantic marker: {pdf_inspection}")
    metrics = image_metrics(WORKSPACE / images[0])
    if metrics["width"] < 600 or metrics["height"] < 400:
        raise AssertionError(f"{kind} preview resolution regressed: {metrics}")
    expected_landscape = kind in {"xlsx", "pptx"}
    if (metrics["width"] > metrics["height"]) != expected_landscape:
        raise AssertionError(f"{kind} preview orientation regressed: {metrics}")
    minimum_non_white = 0.20 if kind == "pptx" else 0.004
    if metrics["non_white_ratio"] < minimum_non_white or metrics["saturated_ratio"] < 0.001:
        raise AssertionError(f"{kind} preview became blank or lost its color styling: {metrics}")
    if kind == "pptx" and metrics["dark_ratio"] < 0.20:
        raise AssertionError(f"PowerPoint dark-background rendering regressed: {metrics}")
    return metrics


WORKSPACE.mkdir(parents=True, exist_ok=True)
fixtures = {
    "docx": WORKSPACE / "visual-baseline.docx",
    "xlsx": WORKSPACE / "visual-baseline.xlsx",
    "pptx": WORKSPACE / "visual-baseline.pptx",
}
create_docx(fixtures["docx"])
create_xlsx(fixtures["xlsx"])
create_pptx(fixtures["pptx"])
results = {kind: verify_render(kind, path) for kind, path in fixtures.items()}
print(json.dumps({"office_visual_regression": "ok", "renders": results}, sort_keys=True))
