#!/usr/bin/env python3
import json
from pathlib import Path
import shutil
import subprocess
import sys
from datetime import datetime, timezone
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

WORKSPACE = Path("/workspace").resolve()
ARTIFACTS = WORKSPACE / "artifacts" / "office"
STAMP = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%S%fZ")
OOXML = {".docx", ".xlsx", ".pptx"}
LEGACY = {".doc", ".xls", ".ppt"}
BLOCKED_ACTIVE = {".docm", ".dotm", ".xlsm", ".xltm", ".xlam", ".pptm", ".potm", ".ppam", ".sldm"}


def workspace_path(value: str, *, must_exist: bool = False) -> Path:
    if not isinstance(value, str) or not value or "\\" in value:
        raise ValueError("path must be a non-empty relative POSIX path")
    candidate = (WORKSPACE / value).resolve()
    if candidate != WORKSPACE and WORKSPACE not in candidate.parents:
        raise ValueError("path must stay inside the run workspace")
    if candidate.suffix.lower() in BLOCKED_ACTIVE:
        raise ValueError("macro-enabled Office formats are blocked by policy")
    if must_exist and not candidate.is_file():
        raise ValueError(f"file does not exist: {value}")
    return candidate


def relative(path: Path) -> str:
    return path.resolve().relative_to(WORKSPACE).as_posix()


def run(command: list[str], timeout: int = 180) -> subprocess.CompletedProcess[str]:
    result = subprocess.run(command, text=True, capture_output=True, timeout=timeout, check=False)
    if result.returncode != 0:
        raise RuntimeError(f"command failed ({result.returncode}): {result.stderr[-4000:]}")
    return result


def libreoffice_convert(source: Path, output_dir: Path, extension: str) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    run([
        "libreoffice", "--headless", "--safe-mode", "--nologo", "--nodefault",
        "--nofirststartwizard", "--nolockcheck", "--invisible", "--convert-to",
        extension, "--outdir", str(output_dir), str(source),
    ])
    output = output_dir / f"{source.stem}.{extension.split(':', 1)[0]}"
    if not output.is_file():
        raise RuntimeError(f"LibreOffice did not create {output.name}")
    return output


def inspect_docx(path: Path) -> dict[str, Any]:
    document = Document(path)
    paragraphs = [item.text for item in document.paragraphs]
    tables = [[[cell.text for cell in row.cells] for row in table.rows] for table in document.tables]
    return {"type": "word", "paragraphs": paragraphs, "tables": tables}


def inspect_xlsx(path: Path) -> dict[str, Any]:
    workbook = load_workbook(path, read_only=True, data_only=False, keep_links=False)
    sheets = []
    for sheet in workbook.worksheets:
        cells = []
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None:
                    cells.append({"cell": cell.coordinate, "value": cell.value, "data_type": cell.data_type})
                if len(cells) >= 20_000:
                    break
            if len(cells) >= 20_000:
                break
        sheets.append({"name": sheet.title, "max_row": sheet.max_row, "max_column": sheet.max_column, "cells": cells})
    workbook.close()
    return {"type": "excel", "sheets": sheets}


def inspect_pptx(path: Path) -> dict[str, Any]:
    presentation = Presentation(path)
    slides = []
    for index, slide in enumerate(presentation.slides, 1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text_frame") and shape.has_text_frame]
        slides.append({"number": index, "texts": texts, "shapes": len(slide.shapes)})
    return {"type": "powerpoint", "slides": slides, "slide_width": presentation.slide_width, "slide_height": presentation.slide_height}


def inspect_pdf(path: Path) -> dict[str, Any]:
    reader = PdfReader(path)
    return {
        "type": "pdf",
        "pages": len(reader.pages),
        "metadata": {str(key): str(value) for key, value in (reader.metadata or {}).items()},
        "text": "\n".join((page.extract_text() or "") for page in reader.pages)[:200_000],
    }


def inspect(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    if suffix in LEGACY:
        converted_extension = {".doc": "docx", ".xls": "xlsx", ".ppt": "pptx"}[suffix]
        path = libreoffice_convert(path, ARTIFACTS / "converted", converted_extension)
        suffix = path.suffix.lower()
    if suffix == ".docx":
        result = inspect_docx(path)
    elif suffix == ".xlsx":
        result = inspect_xlsx(path)
    elif suffix == ".pptx":
        result = inspect_pptx(path)
    elif suffix == ".pdf":
        result = inspect_pdf(path)
    else:
        raise ValueError("supported formats are DOC/DOCX, XLS/XLSX, PPT/PPTX and PDF")
    result["path"] = relative(path)
    return result


def replace_runs(paragraph: Any, old: str, new: str, replace_all: bool) -> int:
    text = "".join(run.text for run in paragraph.runs)
    count = text.count(old)
    if not count:
        return 0
    replacements = count if replace_all else 1
    updated = text.replace(old, new) if replace_all else text.replace(old, new, 1)
    if paragraph.runs:
        paragraph.runs[0].text = updated
        for run_item in paragraph.runs[1:]:
            run_item.text = ""
    else:
        paragraph.text = updated
    return replacements


def replace_text(source: Path, target: Path, old: str, new: str, replace_all: bool) -> int:
    suffix = source.suffix.lower()
    replacements = 0
    if suffix == ".docx":
        document = Document(source)
        for paragraph in document.paragraphs:
            replacements += replace_runs(paragraph, old, new, replace_all)
            if replacements and not replace_all:
                break
        if replace_all or not replacements:
            for table in document.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            replacements += replace_runs(paragraph, old, new, replace_all)
                            if replacements and not replace_all:
                                break
        target.parent.mkdir(parents=True, exist_ok=True)
        document.save(target)
    elif suffix == ".xlsx":
        workbook = load_workbook(source, keep_links=False, keep_vba=False)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and old in cell.value:
                        count = cell.value.count(old)
                        cell.value = cell.value.replace(old, new) if replace_all else cell.value.replace(old, new, 1)
                        replacements += count if replace_all else 1
                        if not replace_all:
                            break
                if replacements and not replace_all:
                    break
            if replacements and not replace_all:
                break
        target.parent.mkdir(parents=True, exist_ok=True)
        workbook.save(target)
    elif suffix == ".pptx":
        presentation = Presentation(source)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    replacements += replace_runs(paragraph, old, new, replace_all)
                    if replacements and not replace_all:
                        break
                if replacements and not replace_all:
                    break
            if replacements and not replace_all:
                break
        target.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(target)
    else:
        raise ValueError("structured replacement requires DOCX, XLSX or PPTX")
    if not replacements:
        target.unlink(missing_ok=True)
        raise ValueError("old_text was not found")
    return replacements


request = json.load(sys.stdin)
action = request.get("action")
source = workspace_path(request.get("path", ""), must_exist=True)
ARTIFACTS.mkdir(parents=True, exist_ok=True)

if action == "inspect":
    response = inspect(source)
elif action == "replace_text":
    target = workspace_path(request.get("output_path") or request["path"])
    replacements = replace_text(source, target, str(request["old_text"]), str(request["new_text"]), bool(request.get("replace_all", False)))
    response = {"path": relative(target), "replacements": replacements, "artifacts": [relative(target)]}
elif action == "export_pdf":
    if source.suffix.lower() == ".pdf":
        target = workspace_path(request.get("output_path") or f"artifacts/office/{source.stem}-{STAMP}.pdf")
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, target)
    else:
        generated = libreoffice_convert(source, ARTIFACTS / ".render" / STAMP, "pdf")
        target = workspace_path(request.get("output_path") or f"artifacts/office/{source.stem}-{source.suffix.lstrip('.')}-{STAMP}.pdf")
        target.parent.mkdir(parents=True, exist_ok=True)
        if generated != target:
            generated.replace(target)
    response = {"path": relative(target), "artifacts": [relative(target)]}
elif action == "preview":
    if source.suffix.lower() == ".pdf":
        pdf = source
    else:
        generated = libreoffice_convert(source, ARTIFACTS / ".render" / STAMP, "pdf")
        pdf = ARTIFACTS / f"{source.stem}-{source.suffix.lstrip('.')}-{STAMP}.pdf"
        generated.replace(pdf)
    prefix = ARTIFACTS / f"{source.stem}-{STAMP}-preview"
    command = ["pdftoppm", "-png", "-r", str(min(int(request.get("dpi", 120)), 200))]
    if not request.get("all_pages", False):
        command += ["-f", "1", "-singlefile"]
    command += [str(pdf), str(prefix)]
    run(command)
    images = sorted(ARTIFACTS.glob(f"{prefix.name}*.png"))
    response = {"pdf": relative(pdf), "images": [relative(path) for path in images], "artifacts": [relative(pdf), *[relative(path) for path in images]]}
else:
    raise ValueError(f"unsupported office action: {action}")

print(json.dumps(response, ensure_ascii=False, default=str))
