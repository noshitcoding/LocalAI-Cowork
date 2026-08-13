#!/usr/bin/env python3
"""Create inert OOXML fixtures without starting Microsoft Office."""

from pathlib import Path
import sys

from docx import Document
from openpyxl import Workbook
from pptx import Presentation


def main() -> None:
    if len(sys.argv) != 2:
        raise SystemExit("usage: create-office-fixtures.py OUTPUT_DIRECTORY")
    root = Path(sys.argv[1]).resolve()
    root.mkdir(parents=True, exist_ok=True)

    word = Document()
    word.add_paragraph("Hello old text")
    word.save(root / "source.docx")

    excel = Workbook()
    sheet = excel.active
    sheet.title = "Data"
    sheet["A1"] = "Category"
    sheet["B1"] = 2
    sheet["A2"] = "Second"
    sheet["B2"] = 3
    excel.save(root / "source.xlsx")

    powerpoint = Presentation()
    slide = powerpoint.slides.add_slide(powerpoint.slide_layouts[1])
    slide.shapes.title.text = "Initial slide"
    slide.placeholders[1].text = "Initial body"
    powerpoint.save(root / "source.pptx")


if __name__ == "__main__":
    main()
