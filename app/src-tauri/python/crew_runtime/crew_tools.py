from __future__ import annotations

import fnmatch
import base64
import html
import ipaddress
import json
import os
import re
import shutil
import socket
import subprocess
import sys
import tempfile
import urllib.parse
import urllib.request
import xml.etree.ElementTree as ET
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Literal

from crewai.tools import BaseTool  # type: ignore
from pydantic import BaseModel, Field, PrivateAttr


MAX_FILE_READ_BYTES = 250_000
MAX_FILE_WRITE_CHARS = 1_500_000
MAX_WEB_BYTES = 1_000_000
MAX_TOOL_OUTPUT_CHARS = 24_000
IGNORED_DIRECTORY_NAMES = {".git", ".next", ".venv", "dist", "node_modules", "target"}


def _workspace_root(request: dict) -> Path:
    configured = str(request.get("cwd") or "").strip()
    candidate = Path(configured).expanduser() if configured else Path.cwd()
    try:
        resolved = candidate.resolve(strict=True)
    except (OSError, RuntimeError):
        resolved = Path.cwd().resolve()
    return resolved if resolved.is_dir() else resolved.parent


def _authorized_roots(request: dict) -> list[tuple[Path, str]]:
    roots: list[tuple[Path, str]] = []
    seen: set[str] = set()

    configured_cwd = str(request.get("cwd") or "").strip()
    candidates: list[dict[str, str]] = []
    if configured_cwd:
        candidates.append({"path": configured_cwd, "kind": "folder", "access": "read_write"})
    candidates.extend(
        entry for entry in (request.get("authorizedPaths") or [])
        if isinstance(entry, dict)
    )

    for entry in candidates:
        raw = str(entry.get("path") or "").strip()
        kind = str(entry.get("kind") or "").strip()
        access = str(entry.get("access") or "read_write").strip()
        if not raw or kind not in {"file", "folder"} or access != "read_write":
            continue
        candidate = Path(raw).expanduser()
        try:
            resolved = candidate.resolve(strict=True)
        except (OSError, RuntimeError):
            continue
        if (kind == "folder" and not resolved.is_dir()) or (kind == "file" and not resolved.is_file()):
            continue
        key = os.path.normcase(str(resolved))
        if key not in seen:
            roots.append((resolved, kind))
            seen.add(key)

    # Backwards compatibility for older saved crew snapshots without authorizedPaths.
    if not roots and "authorizedPaths" not in request:
        root = _workspace_root(request)
        roots.append((root, "folder"))
    return roots


def _primary_workspace_root(roots: list[tuple[Path, str]]) -> Path:
    folder = next((path for path, kind in roots if kind == "folder"), None)
    if folder is not None:
        return folder
    file_root = next((path for path, kind in roots if kind == "file"), None)
    return file_root.parent if file_root is not None else Path.cwd().resolve()


def _resolve_workspace_path(
    roots: list[tuple[Path, str]],
    value: str,
    *,
    allow_root: bool = True,
    tool: str = "",
    deny_rules: list[str] | None = None,
) -> Path:
    if not roots:
        raise PermissionError("No file or folder paths are authorized for this run.")
    root = _primary_workspace_root(roots)
    raw = str(value or "").strip()
    if not raw:
        target = root
    else:
        candidate = Path(raw).expanduser()
        target = (candidate if candidate.is_absolute() else root / candidate).resolve(strict=False)
    authorized = False
    for allowed_path, kind in roots:
        if kind == "file":
            authorized = target == allowed_path
        else:
            try:
                target.relative_to(allowed_path)
                authorized = True
            except ValueError:
                authorized = False
        if authorized:
            break
    if not authorized:
        raise ValueError(
            f"Path is outside the authorized working directory or project paths: {value}"
        )
    rendered_target = str(target)
    for rule in deny_rules or []:
        normalized_rule = str(rule or "").strip()
        if not normalized_rule:
            continue
        if ":" in normalized_rule:
            rule_tool, rule_target = normalized_rule.split(":", 1)
        else:
            rule_tool, rule_target = normalized_rule, "*"
        if (
            fnmatch.fnmatchcase(tool.lower(), rule_tool.lower())
            and fnmatch.fnmatchcase(rendered_target.lower(), rule_target.lower())
        ):
            raise PermissionError(f"Global deny rule blocks {tool} for this path.")
    if not allow_root and any(kind == "folder" and target == path for path, kind in roots):
        raise ValueError("The working-directory root itself cannot be modified.")
    return target


def _display_authorized_path(roots: list[tuple[Path, str]], target: Path) -> str:
    for root, kind in roots:
        if kind == "file" and target == root:
            return target.name
        if kind == "folder":
            try:
                return target.relative_to(root).as_posix()
            except ValueError:
                continue
    return str(target)


def _path_deny_rules(request: dict) -> list[str]:
    governance = request.get("governance") or {}
    if governance.get("policyStrict") is False:
        return []
    return [
        str(rule).strip()
        for rule in governance.get("denyRules") or []
        if str(rule).strip()
    ]


def _truncate(value: object, limit: int = MAX_TOOL_OUTPUT_CHARS) -> str:
    text = str(value or "")
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + f"\n...[truncated after {limit} characters]"


def _safe_result(operation: str, callback) -> str:
    try:
        return _truncate(callback())
    except Exception as exc:
        return f"ERROR ({operation}): {exc.__class__.__name__}: {exc}"


def _agent_access(request: dict, agent_id: str) -> dict:
    governance = request.get("governance") or {}
    for entry in governance.get("agentAccess") or []:
        if isinstance(entry, dict) and str(entry.get("agentId") or "").strip() == agent_id:
            return entry
    return {}


def _canonical_tool_id(value: str) -> str:
    normalized = str(value or "").strip().lower().replace("-", "_")
    aliases = {
        "shell": "bash",
        "bashtool": "bash",
        "read": "read_file",
        "filereadtool": "read_file",
        "edit": "edit_file",
        "write": "edit_file",
        "fileedittool": "edit_file",
        "webfetch": "web_fetch",
        "websearch": "web_search",
        "mcp_call": "mcp",
        "generate_office_workflow": "office_workflow",
        "pptx_template_workflow": "office_workflow",
        "docx_template_workflow": "office_workflow",
    }
    return aliases.get(normalized, normalized)


class ReadFileInput(BaseModel):
    path: str = Field(description="Workspace-relative or absolute path inside the working directory")
    start_line: int = Field(default=1, ge=1, description="First 1-based line to return")
    max_lines: int = Field(default=400, ge=1, le=2000, description="Maximum number of lines")


class ReadFileTool(BaseTool):
    name: str = "read_file"
    description: str = "Read a UTF-8 text file inside the authorized working directory with line numbers."
    args_schema: type[BaseModel] = ReadFileInput
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules

    def _run(self, path: str, start_line: int = 1, max_lines: int = 400) -> str:
        def execute() -> str:
            target = _resolve_workspace_path(
                self._roots, path, tool="read_file", deny_rules=self._deny_rules
            )
            if not target.is_file():
                raise FileNotFoundError(f"File not found: {target}")
            if target.stat().st_size > MAX_FILE_READ_BYTES:
                raise ValueError(f"File exceeds the {MAX_FILE_READ_BYTES}-byte read limit")
            raw = target.read_bytes()
            if b"\x00" in raw:
                raise ValueError("Binary files are not supported by read_file")
            text = raw.decode("utf-8", errors="replace")
            lines = text.splitlines()
            start = max(0, start_line - 1)
            selected = lines[start:start + max_lines]
            rendered = "\n".join(f"{index + start + 1}: {line}" for index, line in enumerate(selected))
            return f"File: {target}\nLines: {start + 1}-{start + len(selected)} of {len(lines)}\n\n{rendered}"

        return _safe_result("read_file", execute)


class EditFileInput(BaseModel):
    path: str = Field(description="File path inside the working directory")
    content: str = Field(default="", description="Complete new file content; use this or old_text/new_text")
    old_text: str = Field(default="", description="Exact existing text to replace")
    new_text: str = Field(default="", description="Replacement text used with old_text")
    replace_all: bool = Field(default=False, description="Replace every occurrence of old_text")


class EditFileTool(BaseTool):
    name: str = "edit_file"
    description: str = "Create or edit a UTF-8 text file atomically. Use content for a full write, or old_text/new_text for a precise replacement."
    args_schema: type[BaseModel] = EditFileInput
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules

    def _run(
        self,
        path: str,
        content: str = "",
        old_text: str = "",
        new_text: str = "",
        replace_all: bool = False,
    ) -> str:
        def execute() -> str:
            target = _resolve_workspace_path(
                self._roots,
                path,
                allow_root=False,
                tool="edit_file",
                deny_rules=self._deny_rules,
            )
            if target.exists() and not target.is_file():
                raise ValueError(f"Target is not a file: {target}")
            if old_text:
                if not target.is_file():
                    raise FileNotFoundError(f"Cannot replace text in missing file: {target}")
                current = target.read_text(encoding="utf-8", errors="strict")
                occurrences = current.count(old_text)
                if occurrences == 0:
                    raise ValueError("old_text was not found; re-read the file before editing")
                if occurrences > 1 and not replace_all:
                    raise ValueError(f"old_text occurs {occurrences} times; provide a unique match or set replace_all")
                updated = current.replace(old_text, new_text, -1 if replace_all else 1)
                change = f"replaced {occurrences if replace_all else 1} occurrence(s)"
            else:
                updated = content
                change = "wrote complete content"
            if len(updated) > MAX_FILE_WRITE_CHARS:
                raise ValueError(f"Content exceeds the {MAX_FILE_WRITE_CHARS}-character write limit")
            target.parent.mkdir(parents=True, exist_ok=True)
            with tempfile.NamedTemporaryFile("w", encoding="utf-8", newline="", dir=target.parent, delete=False) as handle:
                handle.write(updated)
                temporary = Path(handle.name)
            os.replace(temporary, target)
            return f"Updated {target} ({change}, {len(updated)} characters)."

        return _safe_result("edit_file", execute)


class PathInput(BaseModel):
    path: str = Field(description="Path inside the working directory")


class CreateDirectoryTool(BaseTool):
    name: str = "create_directory"
    description: str = "Create a directory, including missing parent directories, inside the working directory."
    args_schema: type[BaseModel] = PathInput
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules

    def _run(self, path: str) -> str:
        return _safe_result("create_directory", lambda: self._create(path))

    def _create(self, path: str) -> str:
        target = _resolve_workspace_path(
            self._roots,
            path,
            allow_root=False,
            tool="create_directory",
            deny_rules=self._deny_rules,
        )
        target.mkdir(parents=True, exist_ok=True)
        return f"Created directory: {target}"


class TransferPathInput(BaseModel):
    source: str = Field(description="Existing source path inside the working directory")
    destination: str = Field(description="Destination path inside the working directory")


class MovePathTool(BaseTool):
    name: str = "move_path"
    description: str = "Move or rename a file or directory within the working directory."
    args_schema: type[BaseModel] = TransferPathInput
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules

    def _run(self, source: str, destination: str) -> str:
        def execute() -> str:
            src = _resolve_workspace_path(
                self._roots, source, allow_root=False, tool="move_path", deny_rules=self._deny_rules
            )
            dst = _resolve_workspace_path(
                self._roots, destination, allow_root=False, tool="move_path", deny_rules=self._deny_rules
            )
            if not src.exists():
                raise FileNotFoundError(f"Source does not exist: {src}")
            if dst.exists():
                raise FileExistsError(f"Destination already exists: {dst}")
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(src), str(dst))
            return f"Moved {src} to {dst}"

        return _safe_result("move_path", execute)


class CopyPathTool(BaseTool):
    name: str = "copy_path"
    description: str = "Copy a file or directory within the working directory without overwriting an existing destination."
    args_schema: type[BaseModel] = TransferPathInput
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules

    def _run(self, source: str, destination: str) -> str:
        def execute() -> str:
            src = _resolve_workspace_path(
                self._roots, source, allow_root=False, tool="copy_path", deny_rules=self._deny_rules
            )
            dst = _resolve_workspace_path(
                self._roots, destination, allow_root=False, tool="copy_path", deny_rules=self._deny_rules
            )
            if not src.exists():
                raise FileNotFoundError(f"Source does not exist: {src}")
            if dst.exists():
                raise FileExistsError(f"Destination already exists: {dst}")
            dst.parent.mkdir(parents=True, exist_ok=True)
            if src.is_dir():
                shutil.copytree(src, dst)
            else:
                shutil.copy2(src, dst)
            return f"Copied {src} to {dst}"

        return _safe_result("copy_path", execute)


class GlobInput(BaseModel):
    pattern: str = Field(description="Glob pattern such as **/*.py")
    path: str = Field(default=".", description="Directory inside the working directory")
    max_results: int = Field(default=200, ge=1, le=1000)


class GlobTool(BaseTool):
    name: str = "glob"
    description: str = "Find workspace files by glob pattern."
    args_schema: type[BaseModel] = GlobInput
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules

    def _run(self, pattern: str, path: str = ".", max_results: int = 200) -> str:
        def execute() -> str:
            base = _resolve_workspace_path(
                self._roots, path, tool="glob", deny_rules=self._deny_rules
            )
            if not base.is_dir():
                raise NotADirectoryError(f"Not a directory: {base}")
            matches: list[str] = []
            for candidate in base.glob(pattern):
                try:
                    candidate = _resolve_workspace_path(
                        self._roots,
                        str(candidate),
                        tool="glob",
                        deny_rules=self._deny_rules,
                    )
                except ValueError:
                    continue
                display_path = _display_authorized_path(self._roots, candidate)
                relative_parts = Path(display_path).parts
                if any(part in IGNORED_DIRECTORY_NAMES for part in relative_parts):
                    continue
                matches.append(display_path + ("/" if candidate.is_dir() else ""))
                if len(matches) >= max_results:
                    break
            return f"Found {len(matches)} path(s):\n" + "\n".join(matches)

        return _safe_result("glob", execute)


class GrepInput(BaseModel):
    pattern: str = Field(description="Regular expression or literal text to search for")
    path: str = Field(default=".", description="File or directory inside the working directory")
    file_pattern: str = Field(default="*", description="Filename glob, for example *.py")
    case_sensitive: bool = False
    max_results: int = Field(default=200, ge=1, le=1000)


class GrepTool(BaseTool):
    name: str = "grep"
    description: str = "Search UTF-8 workspace files and return path, line number, and matching line."
    args_schema: type[BaseModel] = GrepInput
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules

    def _run(
        self,
        pattern: str,
        path: str = ".",
        file_pattern: str = "*",
        case_sensitive: bool = False,
        max_results: int = 200,
    ) -> str:
        def execute() -> str:
            target = _resolve_workspace_path(
                self._roots, path, tool="grep", deny_rules=self._deny_rules
            )
            regex = re.compile(pattern, 0 if case_sensitive else re.IGNORECASE)
            candidates = [target] if target.is_file() else target.rglob("*")
            matches: list[str] = []
            for candidate in candidates:
                if not candidate.is_file() or not fnmatch.fnmatch(candidate.name, file_pattern):
                    continue
                try:
                    candidate = _resolve_workspace_path(
                        self._roots,
                        str(candidate),
                        tool="grep",
                        deny_rules=self._deny_rules,
                    )
                except ValueError:
                    continue
                display_path = _display_authorized_path(self._roots, candidate)
                relative_parts = Path(display_path).parts
                if any(part in IGNORED_DIRECTORY_NAMES for part in relative_parts):
                    continue
                try:
                    if candidate.stat().st_size > MAX_FILE_READ_BYTES:
                        continue
                    raw = candidate.read_bytes()
                    if b"\x00" in raw:
                        continue
                    lines = raw.decode("utf-8", errors="replace").splitlines()
                except OSError:
                    continue
                for line_number, line in enumerate(lines, 1):
                    if regex.search(line):
                        matches.append(f"{display_path}:{line_number}: {_truncate(line, 500)}")
                        if len(matches) >= max_results:
                            return f"Found at least {len(matches)} match(es):\n" + "\n".join(matches)
            return f"Found {len(matches)} match(es):\n" + "\n".join(matches)

        return _safe_result("grep", execute)


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.hidden_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style", "noscript", "svg"}:
            self.hidden_depth += 1
        elif tag in {"br", "p", "div", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript", "svg"} and self.hidden_depth:
            self.hidden_depth -= 1
        elif tag in {"p", "div", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self.hidden_depth:
            self.parts.append(data)

    def text(self) -> str:
        value = html.unescape("".join(self.parts)).replace("\r", "")
        lines = [" ".join(line.split()) for line in value.split("\n")]
        return "\n".join(line for line in lines if line)


def _validate_public_url(value: str) -> str:
    parsed = urllib.parse.urlsplit(str(value or "").strip())
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        raise ValueError("Only absolute http/https URLs are allowed")
    if parsed.username or parsed.password:
        raise ValueError("URLs with embedded credentials are not allowed")
    host = parsed.hostname
    try:
        addresses = {entry[4][0] for entry in socket.getaddrinfo(host, parsed.port or (443 if parsed.scheme == "https" else 80), type=socket.SOCK_STREAM)}
    except OSError as exc:
        raise ValueError(f"Could not resolve host {host}: {exc}") from exc
    if not addresses or any(not ipaddress.ip_address(address).is_global for address in addresses):
        raise ValueError("Private, loopback, link-local, and reserved network destinations are blocked")
    return urllib.parse.urlunsplit(parsed)


class _SafeRedirectHandler(urllib.request.HTTPRedirectHandler):
    def redirect_request(self, req, fp, code, msg, headers, newurl):
        return super().redirect_request(req, fp, code, msg, headers, _validate_public_url(newurl))


def _fetch_public_text(url: str) -> tuple[str, str, int, bool]:
    validated = _validate_public_url(url)
    request = urllib.request.Request(
        validated,
        headers={
            "User-Agent": "LocalAICowork-CrewRuntime/1.0 (+local user initiated research)",
            "Accept": "text/html,application/json,text/plain,application/xml;q=0.9,*/*;q=0.2",
        },
    )
    opener = urllib.request.build_opener(_SafeRedirectHandler())
    with opener.open(request, timeout=20) as response:
        content_type = response.headers.get_content_type().lower()
        allowed = content_type.startswith("text/") or content_type in {"application/json", "application/xml", "application/xhtml+xml"}
        if not allowed:
            raise ValueError(f"Unsupported web content type: {content_type}")
        raw = response.read(MAX_WEB_BYTES + 1)
        truncated = len(raw) > MAX_WEB_BYTES
        if truncated:
            raw = raw[:MAX_WEB_BYTES]
        charset = response.headers.get_content_charset() or "utf-8"
        body = raw.decode(charset, errors="replace")
        final_url = _validate_public_url(response.geturl())
        return final_url, body, int(getattr(response, "status", 200)), truncated


class WebFetchInput(BaseModel):
    url: str = Field(description="Public http/https URL")


class WebFetchTool(BaseTool):
    name: str = "web_fetch"
    description: str = "Fetch readable text from a public web URL. Private and local network destinations are blocked."
    args_schema: type[BaseModel] = WebFetchInput

    def _run(self, url: str) -> str:
        def execute() -> str:
            final_url, body, status, truncated = _fetch_public_text(url)
            extractor = _TextExtractor()
            extractor.feed(body)
            text = extractor.text()
            limit_note = "\nDownload truncated safely after 1000000 bytes." if truncated else ""
            return f"URL: {final_url}\nHTTP: {status}{limit_note}\n\n{_truncate(text, 20_000)}"

        return _safe_result("web_fetch", execute)


class _DuckDuckGoParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.results: list[dict[str, str]] = []
        self.current: dict[str, str] | None = None
        self.capture_title = False
        self.capture_snippet = False

    @staticmethod
    def _classes(attrs: list[tuple[str, str | None]]) -> set[str]:
        value = next((value or "" for key, value in attrs if key == "class"), "")
        return set(value.split())

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        classes = self._classes(attrs)
        if tag == "a" and ("result__a" in classes or "result-link" in classes):
            href = next((value or "" for key, value in attrs if key == "href"), "")
            self.current = {"url": href, "title": "", "snippet": ""}
            self.capture_title = True
        elif self.current is not None and ("result__snippet" in classes or "result-snippet" in classes):
            self.capture_snippet = True

    def handle_endtag(self, tag: str) -> None:
        if tag == "a" and self.capture_title:
            self.capture_title = False
            if self.current is not None:
                self.results.append(self.current)
        if self.capture_snippet and tag in {"a", "div", "td"}:
            self.capture_snippet = False
            self.current = None

    def handle_data(self, data: str) -> None:
        if self.current is None:
            return
        if self.capture_title:
            self.current["title"] += data
        elif self.capture_snippet:
            self.current["snippet"] += data


def _unwrap_search_url(value: str) -> str:
    absolute = urllib.parse.urljoin("https://duckduckgo.com", html.unescape(value))
    parsed = urllib.parse.urlsplit(absolute)
    query = urllib.parse.parse_qs(parsed.query)
    redirected = query.get("uddg", [""])[0]
    return urllib.parse.unquote(redirected) if redirected else absolute


class _BingParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.results: list[dict[str, str]] = []
        self.current: dict[str, str] | None = None
        self.result_depth = 0
        self.heading_depth = 0
        self.capture_title = False
        self.capture_snippet = False

    @staticmethod
    def _classes(attrs: list[tuple[str, str | None]]) -> set[str]:
        value = next((value or "" for key, value in attrs if key == "class"), "")
        return set(value.split())

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        classes = self._classes(attrs)
        if self.current is None and tag == "li" and "b_algo" in classes:
            self.current = {"url": "", "title": "", "snippet": ""}
            self.result_depth = 1
            return
        if self.current is None:
            return
        self.result_depth += 1
        if tag == "h2":
            self.heading_depth = self.result_depth
        elif tag == "a" and self.heading_depth:
            self.current["url"] = next((value or "" for key, value in attrs if key == "href"), "")
            self.capture_title = True
        elif tag == "p" and self.current.get("url"):
            self.capture_snippet = True

    def handle_endtag(self, tag: str) -> None:
        if self.current is None:
            return
        if tag == "a" and self.capture_title:
            self.capture_title = False
        elif tag == "p" and self.capture_snippet:
            self.capture_snippet = False
        elif tag == "h2":
            self.heading_depth = 0
        self.result_depth -= 1
        if self.result_depth <= 0:
            if self.current.get("url") and self.current.get("title"):
                self.results.append(self.current)
            self.current = None
            self.result_depth = 0

    def handle_data(self, data: str) -> None:
        if self.current is None:
            return
        if self.capture_title:
            self.current["title"] += data
        elif self.capture_snippet:
            self.current["snippet"] += data


def _unwrap_bing_url(value: str) -> str:
    absolute = html.unescape(value)
    parsed = urllib.parse.urlsplit(absolute)
    if parsed.hostname and parsed.hostname.lower().endswith("bing.com"):
        encoded = urllib.parse.parse_qs(parsed.query).get("u", [""])[0]
        if encoded.startswith("a1"):
            try:
                payload = encoded[2:]
                payload += "=" * (-len(payload) % 4)
                decoded = base64.urlsafe_b64decode(payload.encode("ascii")).decode("utf-8")
                if decoded.startswith(("http://", "https://")):
                    return decoded
            except Exception:
                pass
    return absolute


class WebSearchInput(BaseModel):
    query: str = Field(description="Specific web search query")
    max_results: int = Field(default=6, ge=1, le=10)


_NEWS_QUERY_PATTERN = re.compile(
    r"(?i)\b(news|nachrichten|breaking|headlines?|stories|latest|aktuell|heute|today)\b"
)
_BROAD_NEWS_QUERY_PATTERN = re.compile(
    r"(?i)\b(world|global|international|welt|worldwide|around\s+the\s+world)\b"
)


def _looks_like_news_query(query: str) -> bool:
    return bool(_NEWS_QUERY_PATTERN.search(query))


def _google_news_feed_url(query: str) -> str:
    locale = "hl=en-US&gl=US&ceid=US:en"
    if _BROAD_NEWS_QUERY_PATTERN.search(query):
        return f"https://news.google.com/rss?{locale}"
    return (
        "https://news.google.com/rss/search?q="
        + urllib.parse.quote_plus(query)
        + f"&{locale}"
    )


def _plain_feed_description(value: str) -> str:
    if not value:
        return ""
    extractor = _TextExtractor()
    extractor.feed(value)
    return " ".join(extractor.text().split())


def _parse_google_news_feed(body: str, max_results: int) -> list[dict[str, str]]:
    try:
        root = ET.fromstring(body)
    except ET.ParseError:
        return []

    results: list[dict[str, str]] = []
    seen: set[str] = set()
    for item in root.findall(".//item"):
        title = " ".join((item.findtext("title") or "").split())
        url = (item.findtext("link") or "").strip()
        if not title or not url.startswith(("http://", "https://")) or url in seen:
            continue

        seen.add(url)
        source = " ".join((item.findtext("source") or "").split())
        published = " ".join((item.findtext("pubDate") or "").split())
        snippet = _plain_feed_description(item.findtext("description") or "")
        results.append({
            "url": url,
            "title": title,
            "snippet": snippet[:700],
            "source": source,
            "published": published,
        })
        if len(results) >= max_results:
            break
    return results


class WebSearchTool(BaseTool):
    name: str = "web_search"
    description: str = (
        "Search the live public web and return titles, source URLs, snippets, and news publication times. "
        "The runtime date is authoritative even when it is newer than the model's training data."
    )
    args_schema: type[BaseModel] = WebSearchInput

    def _run(self, query: str, max_results: int = 6) -> str:
        def execute() -> str:
            normalized = " ".join(str(query or "").split())
            if not normalized:
                raise ValueError("A non-empty search query is required")

            if _looks_like_news_query(normalized):
                provider = "Google News RSS"
                _, body, _, _ = _fetch_public_text(_google_news_feed_url(normalized))
                news_results = _parse_google_news_feed(body, max_results)
                if news_results:
                    rendered = []
                    for index, item in enumerate(news_results, start=1):
                        details = [
                            f"{index}. {item['title']}",
                            f"URL: {item['url']}",
                        ]
                        if item["source"]:
                            details.append(f"Source: {item['source']}")
                        if item["published"]:
                            details.append(f"Published: {item['published']}")
                        details.append(f"Snippet: {item['snippet'] or '(no snippet)'}")
                        rendered.append("\n".join(details))
                    return (
                        f"Search query: {normalized}\nProvider: {provider}\n"
                        f"Results: {len(rendered)}\n\n" + "\n\n".join(rendered)
                    )

            provider = "Bing"
            search_url = (
                "https://www.bing.com/search?q="
                + urllib.parse.quote_plus(normalized)
                + "&setlang=en-US&cc=US&ensearch=1"
            )
            _, body, _, _ = _fetch_public_text(search_url)
            parser: _BingParser | _DuckDuckGoParser = _BingParser()
            parser.feed(body)
            if not parser.results:
                provider = "DuckDuckGo"
                search_url = "https://html.duckduckgo.com/html/?q=" + urllib.parse.quote_plus(normalized)
                _, body, _, _ = _fetch_public_text(search_url)
                parser = _DuckDuckGoParser()
                parser.feed(body)
            rendered: list[str] = []
            seen: set[str] = set()
            for item in parser.results:
                raw_url = item.get("url", "")
                url = _unwrap_bing_url(raw_url) if provider == "Bing" else _unwrap_search_url(raw_url)
                if not url.startswith(("http://", "https://")) or url in seen:
                    continue
                seen.add(url)
                title = " ".join(item.get("title", "").split()) or url
                snippet = " ".join(item.get("snippet", "").split())
                rendered.append(f"{len(rendered) + 1}. {title}\nURL: {url}\nSnippet: {snippet or '(no snippet)'}")
                if len(rendered) >= max_results:
                    break
            if not rendered:
                raise RuntimeError("The search provider returned no parseable results; refine the query and retry")
            return f"Search query: {normalized}\nProvider: {provider}\nResults: {len(rendered)}\n\n" + "\n\n".join(rendered)

        return _safe_result("web_search", execute)


class BashInput(BaseModel):
    command: str = Field(description="Non-interactive command to run in the working directory")
    timeout_seconds: int = Field(default=60, ge=1, le=120)


def _subprocess_environment() -> dict[str, str]:
    blocked_python_variables = {
        "PYTHONHOME",
        "PYTHONPATH",
        "PYTHONEXECUTABLE",
        "__PYVENV_LAUNCHER__",
    }
    environment = {
        key: value
        for key, value in os.environ.items()
        if not re.search(r"(?i)(api[_-]?key|token|secret|password|credential)", key)
        and key.upper() not in blocked_python_variables
    }
    runtime_bin = str(Path(sys.executable).resolve().parent)
    existing_path = environment.get("PATH", "")
    environment["PATH"] = runtime_bin + (os.pathsep + existing_path if existing_path else "")
    return environment


class BashTool(BaseTool):
    name: str = "bash"
    description: str = "Run a bounded, non-interactive PowerShell command on Windows or POSIX shell command elsewhere, from the working directory."
    args_schema: type[BaseModel] = BashInput
    _root: Path = PrivateAttr()
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules
        self._root = _primary_workspace_root(roots)

    def _run(self, command: str, timeout_seconds: int = 60) -> str:
        def execute() -> str:
            normalized = str(command or "").strip()
            if not normalized:
                raise ValueError("A non-empty command is required")
            destructive = re.compile(
                r"(?i)(git\s+(?:reset\s+--hard|clean\s+-[^\n]*f)|remove-item\b[^\n]*(?:-recurse|-force)|\brmdir\s+/s\b|\brd\s+/s\b|\bformat(?:\.com)?\b|\bshutdown\b|\bstop-computer\b)"
            )
            if destructive.search(normalized):
                raise PermissionError("Destructive shell commands are blocked by the Crew runtime")
            if re.search(r"(^|[\\/\s\"'])\.\.([\\/\s\"']|$)", normalized):
                raise PermissionError("Path traversal is blocked by the Crew runtime")
            absolute_candidates = set(
                re.findall(r"(?i)(?:[a-z]:[\\/][^\\s\"'|;&]+|/(?:[^\\s\"'|;&/]+/)*[^\\s\"'|;&]+)", normalized)
            )
            for candidate in absolute_candidates:
                _resolve_workspace_path(
                    self._roots, candidate, tool="bash", deny_rules=self._deny_rules
                )
            relative_candidates = set(
                re.findall(r"(?:^|[\s\"'=])([A-Za-z0-9_.-]+[\\/][^\s\"'|;&]+)", normalized)
            )
            for candidate in relative_candidates:
                candidate_path = self._root / candidate
                if candidate_path.exists():
                    _resolve_workspace_path(
                        self._roots,
                        str(candidate_path),
                        tool="bash",
                        deny_rules=self._deny_rules,
                    )
            args = (
                ["powershell.exe", "-NoProfile", "-NonInteractive", "-Command", normalized]
                if os.name == "nt"
                else ["/bin/sh", "-lc", normalized]
            )
            completed = subprocess.run(
                args,
                cwd=self._root,
                env=_subprocess_environment(),
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=timeout_seconds,
                check=False,
            )
            output = "\n".join(part for part in [completed.stdout.strip(), completed.stderr.strip()] if part)
            return f"Exit code: {completed.returncode}\n{_truncate(output or '(no output)', 20_000)}"

        return _safe_result("bash", execute)


class TodoInput(BaseModel):
    action: Literal["list", "add", "complete", "clear"]
    item: str = Field(default="", description="Todo text for add, or 1-based item number for complete")


class TodoTool(BaseTool):
    name: str = "todo"
    description: str = "Maintain a small in-memory checklist for this agent run."
    args_schema: type[BaseModel] = TodoInput
    _items: list[dict[str, Any]] = PrivateAttr(default_factory=list)

    def _run(self, action: str, item: str = "") -> str:
        def execute() -> str:
            if action == "add":
                text = " ".join(item.split())
                if not text:
                    raise ValueError("Todo text is required")
                self._items.append({"text": text, "done": False})
            elif action == "complete":
                index = int(item) - 1
                if index < 0 or index >= len(self._items):
                    raise ValueError("Todo number is out of range")
                self._items[index]["done"] = True
            elif action == "clear":
                self._items.clear()
            elif action != "list":
                raise ValueError(f"Unknown todo action: {action}")
            if not self._items:
                return "Todo list is empty."
            return "\n".join(f"{index}. [{'x' if value['done'] else ' '}] {value['text']}" for index, value in enumerate(self._items, 1))

        return _safe_result("todo", execute)


class OfficeWorkflowInput(BaseModel):
    output_path: str = Field(description="Output .pptx or .docx path inside the working directory")
    title: str = Field(description="Document or presentation title")
    sections_json: str = Field(
        description='JSON array of sections/slides. Each item may contain "title", "body", and "bullets" (string array).'
    )


def _parse_office_sections(value: str) -> list[dict[str, Any]]:
    parsed = json.loads(value)
    if isinstance(parsed, dict):
        parsed = parsed.get("sections") or parsed.get("slides") or []
    if not isinstance(parsed, list) or not parsed:
        raise ValueError("sections_json must contain a non-empty JSON array")
    sections: list[dict[str, Any]] = []
    for entry in parsed[:40]:
        if not isinstance(entry, dict):
            continue
        bullets = entry.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        sections.append({
            "title": str(entry.get("title") or "Section").strip(),
            "body": str(entry.get("body") or "").strip(),
            "bullets": [str(item).strip() for item in bullets if str(item).strip()][:20],
        })
    if not sections:
        raise ValueError("sections_json did not contain any valid section objects")
    return sections


class OfficeWorkflowTool(BaseTool):
    name: str = "office_workflow"
    description: str = (
        "Create a real PowerPoint (.pptx) or Word (.docx) artifact. Call the tool directly with output_path, title, "
        "and sections_json. sections_json must be a JSON array such as "
        "[{\"title\":\"Evidence\",\"bullets\":[\"Verified fact\"]}]. Do not return a proposed tool call as text."
    )
    args_schema: type[BaseModel] = OfficeWorkflowInput
    _roots: list[tuple[Path, str]] = PrivateAttr()
    _deny_rules: list[str] = PrivateAttr()

    def __init__(self, roots: list[tuple[Path, str]], deny_rules: list[str]) -> None:
        super().__init__()
        self._roots = roots
        self._deny_rules = deny_rules

    def _run(self, output_path: str, title: str, sections_json: str) -> str:
        def execute() -> str:
            target = _resolve_workspace_path(
                self._roots,
                output_path,
                allow_root=False,
                tool="office_workflow",
                deny_rules=self._deny_rules,
            )
            suffix = target.suffix.lower()
            if suffix not in {".pptx", ".docx"}:
                raise ValueError("output_path must end in .pptx or .docx")
            sections = _parse_office_sections(sections_json)
            target.parent.mkdir(parents=True, exist_ok=True)
            if suffix == ".pptx":
                from pptx import Presentation  # type: ignore

                presentation = Presentation()
                title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
                title_slide.shapes.title.text = title.strip() or "Presentation"
                if len(title_slide.placeholders) > 1:
                    title_slide.placeholders[1].text = "Created by LocalAI Cowork CrewAI"
                for section in sections:
                    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                    slide.shapes.title.text = section["title"]
                    frame = slide.placeholders[1].text_frame
                    frame.clear()
                    items = ([section["body"]] if section["body"] else []) + section["bullets"]
                    for index, item in enumerate(items or [""]):
                        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                        paragraph.text = item
                        paragraph.level = 0
                presentation.save(target)
            else:
                from docx import Document  # type: ignore

                document = Document()
                document.add_heading(title.strip() or "Document", level=0)
                for section in sections:
                    document.add_heading(section["title"], level=1)
                    if section["body"]:
                        document.add_paragraph(section["body"])
                    for bullet in section["bullets"]:
                        document.add_paragraph(bullet, style="List Bullet")
                document.save(target)
            return f"Created {target} with {len(sections)} section(s) ({target.stat().st_size} bytes)."

        return _safe_result("office_workflow", execute)


class McpToolInput(BaseModel):
    server_name: str = Field(description="Exact executor-bound MCP server name")
    tool_name: str = Field(description="Tool name exposed by the selected MCP server")
    arguments: dict[str, Any] = Field(default_factory=dict, description="MCP tool arguments")


def _executor_mcp_bindings(request: dict) -> dict[str, dict[str, Any]]:
    bindings: dict[str, dict[str, Any]] = {}
    for value in request.get("executorMcpBindings") or []:
        if not isinstance(value, dict):
            continue
        name = str(value.get("name") or "").strip()
        transport = str(value.get("transport") or "stdio").strip()
        command = value.get("command")
        args = value.get("args", [])
        environment = value.get("environment", {})
        url = value.get("url")
        headers = value.get("headers", {})
        if transport == "stdio" and (
            name
            and isinstance(command, str)
            and command.strip()
            and isinstance(args, list)
            and all(isinstance(argument, str) for argument in args)
            and isinstance(environment, dict)
            and all(isinstance(key, str) and isinstance(secret, str) for key, secret in environment.items())
        ):
            bindings[name] = {
                "name": name,
                "transport": "stdio",
                "command": command,
                "args": args,
                "environment": environment,
            }
        elif transport == "streamable_http" and (
            name
            and isinstance(url, str)
            and url.startswith("https://")
            and isinstance(headers, dict)
            and all(isinstance(key, str) and isinstance(secret, str) for key, secret in headers.items())
        ):
            bindings[name] = {
                "name": name,
                "transport": "streamable_http",
                "url": url,
                "headers": headers,
            }
    return bindings


def _allowed_mcp_server_names(request: dict, agent: dict) -> list[str]:
    requested = [
        str(value).strip()
        for value in agent.get("mcpServerNames") or []
        if str(value).strip()
    ]
    access = _agent_access(request, str(agent.get("id") or "").strip())
    allowed = {
        str(value).strip()
        for value in access.get("allowedMcpServerNames") or []
        if str(value).strip()
    }
    blocked = {
        str(value).strip()
        for value in access.get("blockedMcpServerNames") or []
        if str(value).strip()
    }
    result: list[str] = []
    seen: set[str] = set()
    for name in requested:
        if name in allowed and name not in blocked and name not in seen:
            result.append(name)
            seen.add(name)
    return result


def _redact_mcp_binding_values(value: str, binding: dict[str, Any]) -> str:
    redacted = value
    environment = binding.get("environment") or binding.get("headers") or {}
    secrets = sorted(
        (secret for secret in environment.values() if isinstance(secret, str) and secret),
        key=len,
        reverse=True,
    )
    for secret in secrets:
        redacted = redacted.replace(secret, "[REDACTED]")
        escaped = json.dumps(secret, ensure_ascii=False)[1:-1]
        if escaped != secret:
            redacted = redacted.replace(escaped, "[REDACTED]")
    return redacted


def _redact_mcp_binding_payload(value: Any, binding: dict[str, Any]) -> Any:
    if isinstance(value, str):
        return _redact_mcp_binding_values(value, binding)
    if isinstance(value, list):
        return [_redact_mcp_binding_payload(item, binding) for item in value]
    if isinstance(value, dict):
        return {
            key: _redact_mcp_binding_payload(item, binding)
            for key, item in value.items()
        }
    return value


class McpTool(BaseTool):
    name: str = "mcp_tool"
    description: str = "Call a tool on an executor-bound MCP server."
    args_schema: type[BaseModel] = McpToolInput
    _bindings: dict[str, dict[str, Any]] = PrivateAttr()
    _allowed_names: set[str] = PrivateAttr()
    _root: Path = PrivateAttr()

    def __init__(
        self,
        roots: list[tuple[Path, str]],
        bindings: dict[str, dict[str, Any]],
        allowed_names: list[str],
    ) -> None:
        names = [name for name in allowed_names if name in bindings]
        super().__init__(
            description=(
                "Call a tool on an encrypted executor-bound MCP server. "
                f"Allowed servers: {', '.join(names)}"
            )
        )
        self._bindings = {name: bindings[name] for name in names}
        self._allowed_names = set(names)
        self._root = _primary_workspace_root(roots)

    def _run(
        self,
        server_name: str,
        tool_name: str,
        arguments: dict[str, Any] | None = None,
    ) -> str:
        def execute() -> str:
            name = str(server_name or "").strip()
            tool = str(tool_name or "").strip()
            if name not in self._allowed_names or name not in self._bindings:
                raise PermissionError(f"MCP server is not allowed for this agent: {name}")
            if not tool or len(tool) > 1024 or any(
                ord(character) < 32 or ord(character) == 127 for character in tool
            ):
                raise ValueError("MCP tool name is missing or invalid")
            normalized_arguments = {} if arguments is None else arguments
            if not isinstance(normalized_arguments, dict):
                raise ValueError("MCP tool arguments must be an object")
            payload = {
                "server": self._bindings[name],
                "tool_name": tool,
                "arguments": normalized_arguments,
                "timeout_seconds": 120,
            }
            tool_path = os.environ.get("COWORK_MCP_TOOL_PATH") or "/opt/cowork/mcp-tool.py"
            completed = subprocess.run(
                [sys.executable, tool_path],
                cwd=self._root,
                env=_subprocess_environment(),
                input=json.dumps(payload, ensure_ascii=False),
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=150,
                check=False,
            )
            stdout = _redact_mcp_binding_values(completed.stdout, self._bindings[name])
            stderr = _redact_mcp_binding_values(completed.stderr, self._bindings[name])
            if completed.returncode != 0:
                raise RuntimeError(
                    f"MCP adapter exited with {completed.returncode}: {stderr.strip() or stdout.strip()}"
                )
            response = _redact_mcp_binding_payload(
                json.loads(completed.stdout),
                self._bindings[name],
            )
            if not isinstance(response, dict) or response.get("success") is not True:
                detail = response.get("error") if isinstance(response, dict) else "invalid response"
                raise RuntimeError(f"MCP call failed: {detail}")
            return json.dumps(response.get("result"), ensure_ascii=False, indent=2)

        return _safe_result("mcp_tool", execute)


TOOL_FACTORIES = {
    "read_file": lambda roots, deny_rules: ReadFileTool(roots, deny_rules),
    "edit_file": lambda roots, deny_rules: EditFileTool(roots, deny_rules),
    "create_directory": lambda roots, deny_rules: CreateDirectoryTool(roots, deny_rules),
    "move_path": lambda roots, deny_rules: MovePathTool(roots, deny_rules),
    "copy_path": lambda roots, deny_rules: CopyPathTool(roots, deny_rules),
    "glob": lambda roots, deny_rules: GlobTool(roots, deny_rules),
    "grep": lambda roots, deny_rules: GrepTool(roots, deny_rules),
    "web_fetch": lambda roots, deny_rules: WebFetchTool(),
    "web_search": lambda roots, deny_rules: WebSearchTool(),
    "bash": lambda roots, deny_rules: BashTool(roots, deny_rules),
    "todo": lambda roots, deny_rules: TodoTool(),
    "office_workflow": lambda roots, deny_rules: OfficeWorkflowTool(roots, deny_rules),
}


def build_runtime_tools(request: dict, agent: dict) -> list[BaseTool]:
    agent_id = str(agent.get("id") or "").strip()
    requested = [_canonical_tool_id(value) for value in agent.get("tools") or [] if str(value).strip()]
    access = _agent_access(request, agent_id)
    allowed = {_canonical_tool_id(value) for value in access.get("allowedTools") or []}
    blocked = {_canonical_tool_id(value) for value in access.get("blockedTools") or []}
    roots = _authorized_roots(request)
    deny_rules = _path_deny_rules(request)
    result: list[BaseTool] = []
    seen: set[str] = set()
    for tool_id in requested:
        if tool_id in seen or tool_id in blocked or tool_id not in allowed:
            continue
        if tool_id == "bash" and not any(kind == "folder" for _, kind in roots):
            continue
        factory = TOOL_FACTORIES.get(tool_id)
        if factory is None:
            continue
        result.append(factory(roots, deny_rules))
        seen.add(tool_id)
    bindings = _executor_mcp_bindings(request)
    allowed_mcp_names = _allowed_mcp_server_names(request, agent)
    if any(name in bindings for name in allowed_mcp_names):
        result.append(McpTool(roots, bindings, allowed_mcp_names))
    return result


def unavailable_runtime_tools(request: dict, agent: dict) -> list[str]:
    agent_id = str(agent.get("id") or "").strip()
    requested = {_canonical_tool_id(value) for value in agent.get("tools") or [] if str(value).strip()}
    access = _agent_access(request, agent_id)
    allowed = {_canonical_tool_id(value) for value in access.get("allowedTools") or []}
    unavailable = {
        tool_id
        for tool_id in requested & allowed
        if tool_id not in TOOL_FACTORIES and tool_id not in {"delegate_task", "mcp"}
    }
    bindings = _executor_mcp_bindings(request)
    unavailable.update(
        f"mcp:{name}"
        for name in _allowed_mcp_server_names(request, agent)
        if name not in bindings
    )
    return sorted(unavailable)
